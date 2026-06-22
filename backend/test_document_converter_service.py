import asyncio
from io import BytesIO

from docx import Document
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation
from fastapi import UploadFile
from pypdf import PdfReader
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas

from api.document_converter import DocumentConverterService


async def _upload_file(name: str, content: bytes, content_type: str = "application/octet-stream") -> UploadFile:
    return UploadFile(filename=name, file=BytesIO(content), headers={"content-type": content_type})


def _convert(service: DocumentConverterService, upload: UploadFile, target: str):
    return asyncio.run(service.convert(upload, target))


def _make_docx(text: str) -> bytes:
    buf = BytesIO()
    doc = Document()
    doc.add_heading("测试文档", level=1)
    doc.add_paragraph(text)
    doc.save(buf)
    return buf.getvalue()


def _make_pdf(text: str) -> bytes:
    buf = BytesIO()
    page = canvas.Canvas(buf, pagesize=A4)
    page.drawString(72, 760, text)
    page.save()
    return buf.getvalue()


def _make_png() -> bytes:
    buf = BytesIO()
    image = Image.new("RGB", (320, 180), "#2f80ed")
    image.save(buf, format="PNG")
    return buf.getvalue()


def _make_xlsx() -> bytes:
    buf = BytesIO()
    workbook = Workbook()
    sheet = workbook.active
    assert sheet is not None
    sheet.title = "数据"
    sheet.append(["姓名", "分数"])
    sheet.append(["小巧工具箱", 98])
    workbook.save(buf)
    return buf.getvalue()


def _make_pptx() -> bytes:
    buf = BytesIO()
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[5])
    slide.shapes.title.text = "PPT 标题"
    box = slide.shapes.add_textbox(0, 0, 3000000, 1000000)
    box.text = "PPT 内容转换"
    deck.save(buf)
    return buf.getvalue()


def test_xlsx_and_pptx_convert_to_text_and_pdf():
    service = DocumentConverterService(max_file_size=1024 * 1024)

    xlsx_upload = asyncio.run(_upload_file("table.xlsx", _make_xlsx(), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"))
    xlsx_txt = _convert(service, xlsx_upload, "txt")
    assert xlsx_txt["code"] == 200
    assert "小巧工具箱" in xlsx_txt["content"].decode("utf-8")

    pptx_upload = asyncio.run(_upload_file("slides.pptx", _make_pptx(), "application/vnd.openxmlformats-officedocument.presentationml.presentation"))
    pptx_pdf = _convert(service, pptx_upload, "pdf")
    assert pptx_pdf["code"] == 200
    assert pptx_pdf["filename"] == "slides.pdf"
    assert len(PdfReader(BytesIO(pptx_pdf["content"])).pages) >= 1


def test_pdf_operations_merge_extract_compress_edit_and_remove_watermark():
    service = DocumentConverterService(max_file_size=1024 * 1024)
    pdf1 = _make_pdf("First PDF")
    pdf2 = _make_pdf("Second PDF")

    merged = service.operate_pdf("merge", [{"filename": "a.pdf", "content": pdf1}, {"filename": "b.pdf", "content": pdf2}])
    assert merged["code"] == 200
    assert merged["filename"] == "merged.pdf"
    assert len(PdfReader(BytesIO(merged["content"])).pages) == 2

    extracted = service.operate_pdf("extract", [{"filename": "merged.pdf", "content": merged["content"]}], pages="2")
    assert extracted["code"] == 200
    assert len(PdfReader(BytesIO(extracted["content"])).pages) == 1

    compressed = service.operate_pdf("compress", [{"filename": "a.pdf", "content": pdf1}], compression_level="high")
    assert compressed["code"] == 200
    assert compressed["filename"].endswith("_compressed.pdf")
    assert len(PdfReader(BytesIO(compressed["content"])).pages) == 1

    edited = service.operate_pdf("edit", [{"filename": "a.pdf", "content": pdf1}], text="备注")
    assert edited["code"] == 200
    assert edited["filename"].endswith("_edited.pdf")

    cleaned = service.operate_pdf("remove_watermark", [{"filename": "a.pdf", "content": pdf1}], text="First")
    assert cleaned["code"] == 200
    assert cleaned["filename"].endswith("_clean.pdf")


def test_image_can_convert_to_pdf_and_word():
    service = DocumentConverterService(max_file_size=1024 * 1024)

    image_pdf = asyncio.run(_upload_file("photo.png", _make_png(), "image/png"))
    pdf = _convert(service, image_pdf, "pdf")
    assert pdf["code"] == 200
    assert pdf["filename"] == "photo.pdf"
    assert pdf["media_type"] == "application/pdf"
    assert len(PdfReader(BytesIO(pdf["content"])).pages) == 1

    image_word = asyncio.run(_upload_file("photo.png", _make_png(), "image/png"))
    docx = _convert(service, image_word, "word")
    assert docx["code"] == 200
    assert docx["filename"] == "photo.docx"
    doc = Document(BytesIO(docx["content"]))
    assert doc.inline_shapes


def test_txt_can_convert_to_html_docx_and_pdf():
    service = DocumentConverterService(max_file_size=1024 * 1024)
    upload = asyncio.run(_upload_file("note.txt", "第一行\n第二行".encode("utf-8"), "text/plain"))

    html = _convert(service, upload, "html")
    assert html["code"] == 200
    assert html["filename"] == "note.html"
    assert html["media_type"] == "text/html; charset=utf-8"
    assert "第一行" in html["content"].decode("utf-8")

    upload = asyncio.run(_upload_file("note.txt", "第一行\n第二行".encode("utf-8"), "text/plain"))
    docx = _convert(service, upload, "docx")
    assert docx["code"] == 200
    assert docx["filename"] == "note.docx"
    doc = Document(BytesIO(docx["content"]))
    assert "第一行" in "\n".join(p.text for p in doc.paragraphs)

    upload = asyncio.run(_upload_file("note.txt", "Hello PDF\n中文内容".encode("utf-8"), "text/plain"))
    pdf = _convert(service, upload, "pdf")
    assert pdf["code"] == 200
    assert pdf["filename"] == "note.pdf"
    text = "\n".join(page.extract_text() or "" for page in PdfReader(BytesIO(pdf["content"])).pages)
    assert "Hello PDF" in text


def test_docx_pdf_and_html_convert_to_text():
    service = DocumentConverterService(max_file_size=1024 * 1024)

    docx_upload = asyncio.run(_upload_file("demo.docx", _make_docx("Word 内容转换")))
    docx_text = _convert(service, docx_upload, "txt")
    assert docx_text["code"] == 200
    assert docx_text["filename"] == "demo.txt"
    assert "Word 内容转换" in docx_text["content"].decode("utf-8")

    pdf_upload = asyncio.run(_upload_file("demo.pdf", _make_pdf("PDF extract text")))
    pdf_text = _convert(service, pdf_upload, "txt")
    assert pdf_text["code"] == 200
    assert pdf_text["filename"] == "demo.txt"
    assert "PDF extract text" in pdf_text["content"].decode("utf-8")

    html_upload = asyncio.run(_upload_file("demo.html", b"<h1>HTML Title</h1><p>Hello&nbsp;World</p>"))
    html_text = _convert(service, html_upload, "txt")
    assert html_text["code"] == 200
    assert "HTML Title" in html_text["content"].decode("utf-8")
    assert "Hello" in html_text["content"].decode("utf-8")


def test_converter_rejects_unsupported_format_and_too_large_file():
    service = DocumentConverterService(max_file_size=8)

    bad_ext = asyncio.run(_upload_file("demo.exe", b"hello"))
    result = _convert(service, bad_ext, "txt")
    assert result["code"] == 400
    assert "支持" in result["msg"]

    too_large = asyncio.run(_upload_file("demo.txt", b"0123456789"))
    result = _convert(service, too_large, "pdf")
    assert result["code"] == 400
    assert "大小" in result["msg"]
    bad_image_target = asyncio.run(_upload_file("photo.jpg", _make_png(), "image/jpeg"))
    result = _convert(service, bad_image_target, "txt")
    assert result["code"] == 400
    assert "图片" in result["msg"]


def test_same_format_is_rejected_to_avoid_empty_work():
    service = DocumentConverterService(max_file_size=1024)
    upload = asyncio.run(_upload_file("same.txt", b"same"))
    result = _convert(service, upload, "txt")
    assert result["code"] == 400
    assert "请选择不同" in result["msg"]
