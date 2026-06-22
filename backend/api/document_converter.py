import html
import re
from html.parser import HTMLParser
from io import BytesIO
from pathlib import Path
from typing import Any, Dict, List, TypedDict

from docx import Document
from docx.shared import Inches
from fastapi import UploadFile
from openpyxl import Workbook, load_workbook
from PIL import Image, ImageOps
from pptx import Presentation
from pypdf import PdfReader, PdfWriter
from pypdf.generic import ByteStringObject, ContentStream, NameObject, TextStringObject
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.utils import ImageReader
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.pdfgen import canvas
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer


IMAGE_INPUTS = {"jpg", "jpeg", "png", "webp"}
OFFICE_INPUTS = {"xlsx", "xls", "pptx", "ppt"}
SUPPORTED_INPUTS = {"txt", "html", "htm", "docx", "pdf", *IMAGE_INPUTS, *OFFICE_INPUTS}
SUPPORTED_TARGETS = {"txt", "html", "docx", "pdf", "xlsx", "pptx"}
TARGET_ALIASES = {"word": "docx", "text": "txt", "htm": "html", "excel": "xlsx", "ppt": "pptx"}
OUTPUT_MIME_TYPES = {
    "txt": "text/plain; charset=utf-8",
    "html": "text/html; charset=utf-8",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pdf": "application/pdf",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


class PdfInputFile(TypedDict):
    filename: str
    content: bytes


class ScanImageFile(TypedDict):
    filename: str
    content: bytes


class _TextExtractor(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: List[str] = []

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        if tag.lower() in {"p", "br", "div", "li", "h1", "h2", "h3", "tr"}:
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        if data:
            self.parts.append(data)

    def get_text(self) -> str:
        text = html.unescape("".join(self.parts))
        lines = [re.sub(r"\s+", " ", line).strip() for line in text.splitlines()]
        return "\n".join(line for line in lines if line)


class DocumentConverterService:
    """Lightweight document conversion service for common text-oriented formats.

    Supported conversions are text-preserving. PDF conversion extracts selectable
    text from PDF pages; scanned image PDFs need OCR and are intentionally out of
    scope for this lightweight toolbox feature.
    """

    def __init__(self, max_file_size: int = 5 * 1024 * 1024) -> None:
        self.max_file_size = max_file_size

    async def convert(self, upload: UploadFile, target_format: str) -> Dict[str, Any]:
        source_name = Path(upload.filename or "document")
        source_ext = source_name.suffix.lower().lstrip(".")
        source_ext = "html" if source_ext == "htm" else source_ext
        target = TARGET_ALIASES.get(str(target_format or "").lower().strip(), str(target_format or "").lower().strip())

        if source_ext not in SUPPORTED_INPUTS:
            return {"code": 400, "msg": "暂时支持 TXT、HTML、DOCX、PDF、JPG、PNG、WEBP、Excel、PPT 文件"}
        if target not in SUPPORTED_TARGETS:
            return {"code": 400, "msg": "目标格式支持 txt、html、docx、pdf、xlsx、pptx"}
        if source_ext == target:
            return {"code": 400, "msg": "请选择不同的目标格式"}
        if source_ext in IMAGE_INPUTS and target not in {"pdf", "docx"}:
            return {"code": 400, "msg": "图片转换暂时支持转 PDF 或 Word"}
        if source_ext in {"xls", "ppt"}:
            return {"code": 400, "msg": "当前轻量版暂不支持旧版 xls/ppt，请另存为 xlsx/pptx 后再转换"}

        content = await upload.read()
        if not content:
            return {"code": 400, "msg": "文件内容为空"}
        if len(content) > self.max_file_size:
            return {"code": 400, "msg": f"文件大小不能超过 {self.max_file_size // 1024 // 1024}MB"}

        try:
            if source_ext in IMAGE_INPUTS:
                output = self._render_image(content, target, source_name.stem or "image")
            else:
                text = self._extract_text(content, source_ext)
                if not text.strip():
                    return {"code": 400, "msg": "未能从文档中提取到文本内容，扫描版 PDF 暂不支持"}
                output = self._render(text, target, source_name.stem or "document")
        except Exception as exc:  # pragma: no cover - defensive API guard
            return {"code": 400, "msg": f"文档转换失败：{exc}"}

        filename = f"{self._safe_stem(source_name.stem)}.{target}"
        return {
            "code": 200,
            "msg": "success",
            "filename": filename,
            "media_type": OUTPUT_MIME_TYPES[target],
            "content": output,
        }

    def _extract_text(self, content: bytes, source_ext: str) -> str:
        if source_ext == "txt":
            return self._decode_text(content)
        if source_ext == "html":
            parser = _TextExtractor()
            parser.feed(self._decode_text(content))
            return parser.get_text()
        if source_ext == "docx":
            doc = Document(BytesIO(content))
            lines: List[str] = []
            lines.extend(p.text for p in doc.paragraphs if p.text.strip())
            for table in doc.tables:
                for row in table.rows:
                    row_text = "\t".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        lines.append(row_text)
            return "\n".join(lines)
        if source_ext == "pdf":
            reader = PdfReader(BytesIO(content))
            return "\n".join(page.extract_text() or "" for page in reader.pages)
        if source_ext == "xlsx":
            return self._extract_xlsx_text(content)
        if source_ext == "pptx":
            return self._extract_pptx_text(content)
        raise ValueError("不支持的源文件格式")

    def _render(self, text: str, target: str, title: str) -> bytes:
        if target == "txt":
            return text.encode("utf-8")
        if target == "html":
            body = "\n".join(f"<p>{html.escape(line)}</p>" for line in text.splitlines() if line.strip())
            return (
                "<!doctype html><html><head><meta charset=\"utf-8\">"
                f"<title>{html.escape(title)}</title></head><body>{body}</body></html>"
            ).encode("utf-8")
        if target == "docx":
            buf = BytesIO()
            doc = Document()
            for line in text.splitlines():
                doc.add_paragraph(line)
            doc.save(buf)
            return buf.getvalue()
        if target == "pdf":
            return self._text_to_pdf(text, title)
        if target == "xlsx":
            return self._text_to_xlsx(text, title)
        if target == "pptx":
            return self._text_to_pptx(text, title)
        raise ValueError("不支持的目标格式")

    @staticmethod
    def _extract_xlsx_text(content: bytes) -> str:
        workbook = load_workbook(BytesIO(content), data_only=True, read_only=True)
        lines: List[str] = []
        for sheet in workbook.worksheets:
            lines.append(f"[{sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                values = [str(value) for value in row if value is not None and str(value).strip()]
                if values:
                    lines.append("\t".join(values))
        return "\n".join(lines)

    @staticmethod
    def _extract_pptx_text(content: bytes) -> str:
        presentation = Presentation(BytesIO(content))
        lines: List[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            lines.append(f"[Slide {index}]")
            for shape in slide.shapes:
                if hasattr(shape, "text") and str(shape.text).strip():
                    lines.append(str(shape.text).strip())
        return "\n".join(lines)

    def _text_to_pdf(self, text: str, title: str) -> bytes:
        buf = BytesIO()
        doc = SimpleDocTemplate(buf, pagesize=A4, title=title)
        styles = getSampleStyleSheet()
        try:
            pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
            body_style = ParagraphStyle("ToolboxBody", parent=styles["BodyText"], fontName="STSong-Light", fontSize=10, leading=15)
        except Exception:
            body_style = styles["BodyText"]
        story = []
        for line in text.splitlines():
            escaped = html.escape(line) or " "
            story.append(Paragraph(escaped, body_style))
            story.append(Spacer(1, 6))
        doc.build(story or [Paragraph(" ", body_style)])
        return buf.getvalue()

    @staticmethod
    def _text_to_xlsx(text: str, title: str) -> bytes:
        buf = BytesIO()
        workbook = Workbook()
        sheet = workbook.active
        assert sheet is not None
        sheet.title = (title or "Sheet1")[:31]
        for row_index, line in enumerate(text.splitlines() or [""], start=1):
            columns = line.split("\t") if "\t" in line else [line]
            for col_index, value in enumerate(columns, start=1):
                sheet.cell(row=row_index, column=col_index, value=value)
        workbook.save(buf)
        return buf.getvalue()

    @staticmethod
    def _text_to_pptx(text: str, title: str) -> bytes:
        buf = BytesIO()
        deck = Presentation()
        lines = [line for line in text.splitlines() if line.strip()] or [title or "文档内容"]
        chunk_size = 8
        for start in range(0, len(lines), chunk_size):
            slide = deck.slides.add_slide(deck.slide_layouts[1])
            slide.shapes.title.text = title or "文档转换"
            body = slide.placeholders[1].text_frame
            body.clear()
            for line in lines[start:start + chunk_size]:
                paragraph = body.add_paragraph()
                paragraph.text = line[:180]
                paragraph.level = 0
        deck.save(buf)
        return buf.getvalue()

    def _render_image(self, content: bytes, target: str, title: str) -> bytes:
        image = self._open_image(content)
        if target == "pdf":
            return self._image_to_pdf(image, title)
        if target == "docx":
            return self._image_to_docx(image, title)
        raise ValueError("图片转换暂时支持转 PDF 或 Word")

    @staticmethod
    def _open_image(content: bytes) -> Image.Image:
        image = Image.open(BytesIO(content))
        image = ImageOps.exif_transpose(image)
        if image.mode in {"RGBA", "LA", "P"}:
            bg = Image.new("RGB", image.size, "white")
            if image.mode == "P":
                image = image.convert("RGBA")
            bg.paste(image, mask=image.split()[-1] if image.mode in {"RGBA", "LA"} else None)
            return bg
        return image.convert("RGB")

    def _image_to_pdf(self, image: Image.Image, title: str) -> bytes:
        buf = BytesIO()
        page_width, page_height = A4
        margin = 36
        draw_width = page_width - margin * 2
        draw_height = page_height - margin * 2
        ratio = min(draw_width / image.width, draw_height / image.height)
        width = image.width * ratio
        height = image.height * ratio
        x = (page_width - width) / 2
        y = (page_height - height) / 2
        from reportlab.pdfgen import canvas

        page = canvas.Canvas(buf, pagesize=A4)
        page.setTitle(title)
        image_buffer = BytesIO()
        image.save(image_buffer, format="JPEG", quality=92)
        image_buffer.seek(0)
        page.drawImage(ImageReader(image_buffer), x, y, width=width, height=height, preserveAspectRatio=True, mask="auto")
        page.showPage()
        page.save()
        return buf.getvalue()

    @staticmethod
    def _image_to_docx(image: Image.Image, title: str) -> bytes:
        buf = BytesIO()
        doc = Document()
        doc.add_heading(title, level=1)
        image_buffer = BytesIO()
        image.save(image_buffer, format="JPEG", quality=92)
        image_buffer.seek(0)
        doc.add_picture(image_buffer, width=Inches(6))
        doc.save(buf)
        return buf.getvalue()

    def scan_images(self, files: List[ScanImageFile], target_format: str = "pdf", title: str = "扫描文档") -> Dict[str, Any]:
        target = TARGET_ALIASES.get(str(target_format or "").lower().strip(), str(target_format or "").lower().strip())
        if target not in {"pdf", "docx", "pptx"}:
            return {"code": 400, "msg": "扫描生成暂时支持 PDF、Word、PPT"}
        try:
            images = self._validate_scan_images(files)
            if target == "pdf":
                content = self._scan_images_to_pdf(images, title)
            elif target == "docx":
                content = self._scan_images_to_docx(images, title)
            else:
                content = self._scan_images_to_pptx(images, title)
        except Exception as exc:  # pragma: no cover - defensive API guard
            return {"code": 400, "msg": f"扫描生成失败：{exc}"}
        safe_title = self._safe_stem(title or "scan")
        return {
            "code": 200,
            "msg": "success",
            "filename": f"{safe_title}.{target}",
            "media_type": OUTPUT_MIME_TYPES[target],
            "content": content,
        }

    def _validate_scan_images(self, files: List[ScanImageFile]) -> list[Image.Image]:
        if not files:
            raise ValueError("请先拍照或选择图片")
        if len(files) > 10:
            raise ValueError("一次最多支持 10 张扫描图片")
        images: list[Image.Image] = []
        total = 0
        for item in files:
            filename = str(item.get("filename") or "scan.jpg")
            ext = Path(filename).suffix.lower().lstrip(".")
            if ext not in IMAGE_INPUTS:
                raise ValueError("扫描图片支持 JPG、PNG、WEBP")
            content = item.get("content") or b""
            if not content:
                raise ValueError("图片内容为空")
            total += len(content)
            if len(content) > self.max_file_size or total > self.max_file_size * 4:
                raise ValueError(f"单张图片不能超过 {self.max_file_size // 1024 // 1024}MB，总大小不能超过 {self.max_file_size * 4 // 1024 // 1024}MB")
            images.append(self._open_image(content))
        return images

    def _scan_images_to_pdf(self, images: list[Image.Image], title: str) -> bytes:
        buf = BytesIO()
        page_width, page_height = A4
        margin = 28
        draw_width = page_width - margin * 2
        draw_height = page_height - margin * 2
        page = canvas.Canvas(buf, pagesize=A4)
        page.setTitle(title or "扫描文档")
        for image in images:
            ratio = min(draw_width / image.width, draw_height / image.height)
            width = image.width * ratio
            height = image.height * ratio
            x = (page_width - width) / 2
            y = (page_height - height) / 2
            image_buffer = BytesIO()
            image.save(image_buffer, format="JPEG", quality=90)
            image_buffer.seek(0)
            page.drawImage(ImageReader(image_buffer), x, y, width=width, height=height, preserveAspectRatio=True, mask="auto")
            page.showPage()
        page.save()
        return buf.getvalue()

    @staticmethod
    def _scan_images_to_docx(images: list[Image.Image], title: str) -> bytes:
        buf = BytesIO()
        doc = Document()
        doc.add_heading(title or "扫描文档", level=1)
        for index, image in enumerate(images, start=1):
            if index > 1:
                doc.add_page_break()
            image_buffer = BytesIO()
            image.save(image_buffer, format="JPEG", quality=90)
            image_buffer.seek(0)
            doc.add_paragraph(f"扫描页 {index}")
            doc.add_picture(image_buffer, width=Inches(6))
        doc.save(buf)
        return buf.getvalue()

    @staticmethod
    def _scan_images_to_pptx(images: list[Image.Image], title: str) -> bytes:
        buf = BytesIO()
        deck = Presentation()
        blank_layout = deck.slide_layouts[6]
        slide_width = deck.slide_width
        slide_height = deck.slide_height
        for image in images:
            slide = deck.slides.add_slide(blank_layout)
            image_buffer = BytesIO()
            image.save(image_buffer, format="JPEG", quality=90)
            image_buffer.seek(0)
            ratio = min(slide_width / image.width, slide_height / image.height)
            width = int(image.width * ratio)
            height = int(image.height * ratio)
            left = int((slide_width - width) / 2)
            top = int((slide_height - height) / 2)
            slide.shapes.add_picture(image_buffer, left, top, width=width, height=height)
        deck.core_properties.title = title or "扫描文档"
        deck.save(buf)
        return buf.getvalue()

    def operate_pdf(self, operation: str, files: List[PdfInputFile], pages: str = "", text: str = "", compression_level: str = "medium") -> Dict[str, Any]:
        op = str(operation or "").strip().lower()
        try:
            if op == "merge":
                return self._pdf_merge(files)
            if op in {"split", "extract"}:
                return self._pdf_extract(files, pages)
            if op == "compress":
                return self._pdf_compress(files, compression_level)
            if op == "edit":
                return self._pdf_add_text(files, text)
            if op in {"remove_watermark", "watermark_remove"}:
                return self._pdf_remove_watermark(files, text)
        except Exception as exc:  # pragma: no cover - defensive API guard
            return {"code": 400, "msg": f"PDF 处理失败：{exc}"}
        return {"code": 400, "msg": "不支持的 PDF 操作"}

    def _validate_pdf_files(self, files: List[PdfInputFile], min_count: int = 1) -> list[tuple[str, bytes]]:
        if len(files) < min_count:
            return []
        normalized: list[tuple[str, bytes]] = []
        total = 0
        for item in files:
            filename = str(item.get("filename") or "document.pdf")
            content = item.get("content") or b""
            if not filename.lower().endswith(".pdf"):
                raise ValueError("仅支持 PDF 文件")
            if not content:
                raise ValueError("PDF 文件内容为空")
            total += len(content)
            if len(content) > self.max_file_size or total > self.max_file_size * 3:
                raise ValueError(f"单文件不能超过 {self.max_file_size // 1024 // 1024}MB，合并总大小不能超过 {self.max_file_size * 3 // 1024 // 1024}MB")
            normalized.append((filename, content))
        return normalized

    def _pdf_merge(self, files: List[PdfInputFile]) -> Dict[str, Any]:
        pdfs = self._validate_pdf_files(files, min_count=2)
        if not pdfs:
            return {"code": 400, "msg": "PDF 合并至少需要 2 个 PDF 文件"}
        writer = PdfWriter()
        for _, content in pdfs:
            reader = PdfReader(BytesIO(content))
            for page in reader.pages:
                writer.add_page(page)
        return self._pdf_result(writer, "merged.pdf")

    def _pdf_extract(self, files: List[PdfInputFile], pages: str) -> Dict[str, Any]:
        pdfs = self._validate_pdf_files(files, min_count=1)
        if not pdfs:
            return {"code": 400, "msg": "请上传 PDF 文件"}
        filename, content = pdfs[0]
        reader = PdfReader(BytesIO(content))
        indexes = self._parse_page_ranges(pages, len(reader.pages))
        if not indexes:
            return {"code": 400, "msg": "请输入要拆分/提取的页码，例如 1,3-5"}
        writer = PdfWriter()
        for index in indexes:
            writer.add_page(reader.pages[index])
        return self._pdf_result(writer, f"{self._safe_stem(Path(filename).stem)}_pages.pdf")

    def _pdf_compress(self, files: List[PdfInputFile], compression_level: str = "medium") -> Dict[str, Any]:
        pdfs = self._validate_pdf_files(files, min_count=1)
        if not pdfs:
            return {"code": 400, "msg": "请上传 PDF 文件"}
        level = str(compression_level or "medium").strip().lower()
        if level not in {"low", "medium", "high"}:
            level = "medium"
        filename, content = pdfs[0]
        reader = PdfReader(BytesIO(content))
        writer = PdfWriter()
        for page in reader.pages:
            try:
                page.compress_content_streams()
            except Exception:
                pass
            writer.add_page(page)
        if level in {"medium", "high"}:
            try:
                writer.compress_identical_objects(remove_duplicates=True, remove_unreferenced=True)
            except Exception:
                pass
        if level != "high":
            for key, value in (reader.metadata or {}).items():
                if value:
                    try:
                        writer.add_metadata({key: str(value)})
                    except Exception:
                        pass
        return self._pdf_result(writer, f"{self._safe_stem(Path(filename).stem)}_compressed.pdf")

    def _pdf_add_text(self, files: List[PdfInputFile], text: str) -> Dict[str, Any]:
        pdfs = self._validate_pdf_files(files, min_count=1)
        if not pdfs:
            return {"code": 400, "msg": "请上传 PDF 文件"}
        if not str(text or "").strip():
            return {"code": 400, "msg": "请输入要添加到 PDF 的文字"}
        filename, content = pdfs[0]
        reader = PdfReader(BytesIO(content))
        writer = PdfWriter()
        for page in reader.pages:
            width = float(page.mediabox.width)
            height = float(page.mediabox.height)
            overlay = self._make_text_overlay(str(text), width, height)
            page.merge_page(PdfReader(BytesIO(overlay)).pages[0])
            writer.add_page(page)
        return self._pdf_result(writer, f"{self._safe_stem(Path(filename).stem)}_edited.pdf")

    def _pdf_remove_watermark(self, files: List[PdfInputFile], text: str = "") -> Dict[str, Any]:
        pdfs = self._validate_pdf_files(files, min_count=1)
        if not pdfs:
            return {"code": 400, "msg": "请上传 PDF 文件"}
        filename, content = pdfs[0]
        reader = PdfReader(BytesIO(content))
        writer = PdfWriter()
        for page in reader.pages:
            if "/Annots" in page:
                del page[NameObject("/Annots")]
            if str(text or "").strip():
                self._remove_text_from_page_stream(page, str(text).strip())
            try:
                page.compress_content_streams()
            except Exception:
                pass
            writer.add_page(page)
        return self._pdf_result(writer, f"{self._safe_stem(Path(filename).stem)}_clean.pdf")

    @staticmethod
    def _pdf_result(writer: PdfWriter, filename: str) -> Dict[str, Any]:
        buf = BytesIO()
        writer.write(buf)
        return {"code": 200, "msg": "success", "filename": filename, "media_type": OUTPUT_MIME_TYPES["pdf"], "content": buf.getvalue()}

    @staticmethod
    def _parse_page_ranges(spec: str, total_pages: int) -> List[int]:
        indexes: list[int] = []
        for part in str(spec or "").replace("，", ",").split(","):
            part = part.strip()
            if not part:
                continue
            if "-" in part:
                start_s, end_s = part.split("-", 1)
                start, end = int(start_s), int(end_s)
                if start > end:
                    start, end = end, start
                indexes.extend(range(start - 1, end))
            else:
                indexes.append(int(part) - 1)
        seen = []
        for index in indexes:
            if 0 <= index < total_pages and index not in seen:
                seen.append(index)
        return seen

    @staticmethod
    def _make_text_overlay(text: str, width: float, height: float) -> bytes:
        buf = BytesIO()
        page = canvas.Canvas(buf, pagesize=(width, height))
        try:
            pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
            page.setFont("STSong-Light", 14)
        except Exception:
            page.setFont("Helvetica", 14)
        page.setFillColorRGB(0.2, 0.2, 0.2)
        page.drawString(36, max(36, height - 54), text[:160])
        page.save()
        return buf.getvalue()

    @staticmethod
    def _remove_text_from_page_stream(page: Any, watermark_text: str) -> None:
        try:
            content = ContentStream(page.get_contents(), page.pdf)
        except Exception:
            return
        changed = False
        for operands, operator in content.operations:
            if operator in {b"Tj", b"'", b'"'} and operands:
                if watermark_text in str(operands[0]):
                    operands[0] = TextStringObject("") if isinstance(operands[0], TextStringObject) else ByteStringObject(b"")
                    changed = True
            elif operator == b"TJ" and operands:
                items = operands[0]
                for i, item in enumerate(items):
                    if watermark_text in str(item):
                        items[i] = TextStringObject("") if isinstance(item, TextStringObject) else ByteStringObject(b"")
                        changed = True
        if changed:
            page[NameObject("/Contents")] = content

    @staticmethod
    def _decode_text(content: bytes) -> str:
        for encoding in ("utf-8", "utf-8-sig", "gb18030"):
            try:
                return content.decode(encoding)
            except UnicodeDecodeError:
                continue
        return content.decode("utf-8", errors="ignore")

    @staticmethod
    def _safe_stem(stem: str) -> str:
        cleaned = re.sub(r"[^\w\u4e00-\u9fff.-]+", "_", stem or "document").strip("._")
        return cleaned or "document"


_document_converter_service = DocumentConverterService()


def get_document_converter_service() -> DocumentConverterService:
    return _document_converter_service
